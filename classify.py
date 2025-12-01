"""
학습된 KoBERT 모델로 문서를 종류별로 분류하는 스크립트

- input 폴더에 문서 파일들을 넣으면
- output 폴더에 문서 종류(기안서, 견적서, 계약서 등)별로 폴더를 만들어 분류

사용 예시:
python classify.py \
  --model_dir ./saved_kobert_doc_classifier \
  --input_dir ./input \
  --output_dir ./output

지원 문서 형식: .txt, .pdf, .docx, .hwp, .pptx, .xlsx
"""

import os
import shutil
import argparse
from pathlib import Path
from typing import Optional

import torch
import torch.nn as nn
from transformers import BertModel, AutoTokenizer
import json


# ======================
# KoBERT 분류 모델 (train.py와 동일)
# ======================

class KoBERTClassifier(nn.Module):
    def __init__(self, num_classes: int, dropout_rate: float = 0.1):
        super().__init__()
        self.bert = BertModel.from_pretrained("skt/kobert-base-v1")
        self.dropout = nn.Dropout(dropout_rate)
        self.classifier = nn.Linear(768, num_classes)

    def forward(self, input_ids, attention_mask):
        outputs = self.bert(
            input_ids=input_ids,
            attention_mask=attention_mask,
        )
        cls_output = outputs.last_hidden_state[:, 0, :]
        x = self.dropout(cls_output)
        logits = self.classifier(x)
        return logits


# ======================
# 문서에서 텍스트 추출
# ======================

def extract_text_from_txt(filepath: Path) -> str:
    """TXT 파일에서 텍스트 추출"""
    encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
    for enc in encodings:
        try:
            with open(filepath, 'r', encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    return ""


def extract_text_from_pdf(filepath: Path) -> str:
    """PDF 파일에서 텍스트 추출"""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n".join(text_parts)
    except ImportError:
        print("[경고] pdfplumber 설치 필요: pip install pdfplumber")
        return ""
    except Exception as e:
        print(f"[경고] PDF 읽기 실패: {filepath} ({e})")
        return ""


def extract_text_from_docx(filepath: Path) -> str:
    """DOCX 파일에서 텍스트 추출"""
    try:
        from docx import Document
        doc = Document(filepath)
        text_parts = [para.text for para in doc.paragraphs]
        return "\n".join(text_parts)
    except ImportError:
        print("[경고] python-docx 설치 필요: pip install python-docx")
        return ""
    except Exception as e:
        print(f"[경고] DOCX 읽기 실패: {filepath} ({e})")
        return ""


def extract_text_from_hwp(filepath: Path) -> str:
    """HWP 파일에서 텍스트 추출"""
    try:
        import olefile
        with olefile.OleFileIO(filepath) as ole:
            if ole.exists("PrvText"):
                encoded_text = ole.openstream("PrvText").read()
                return encoded_text.decode('utf-16', errors='ignore')
        return ""
    except ImportError:
        print("[경고] olefile 설치 필요: pip install olefile")
        return ""
    except Exception as e:
        print(f"[경고] HWP 읽기 실패: {filepath} ({e})")
        return ""


def extract_text_from_pptx(filepath: Path) -> str:
    """PPTX 파일에서 텍스트 추출"""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except ImportError:
        print("[경고] python-pptx 설치 필요: pip install python-pptx")
        return ""
    except Exception as e:
        print(f"[경고] PPTX 읽기 실패: {filepath} ({e})")
        return ""


def extract_text_from_xlsx(filepath: Path) -> str:
    """XLSX 파일에서 텍스트 추출"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text_parts.append(str(cell.value))
        return " ".join(text_parts)
    except ImportError:
        print("[경고] openpyxl 설치 필요: pip install openpyxl")
        return ""
    except Exception as e:
        print(f"[경고] XLSX 읽기 실패: {filepath} ({e})")
        return ""


def extract_text(filepath: Path) -> str:
    """파일 확장자에 따라 텍스트 추출"""
    ext = filepath.suffix.lower()
    
    extractors = {
        ".txt": extract_text_from_txt,
        ".pdf": extract_text_from_pdf,
        ".doc": extract_text_from_docx,  # .doc은 제한적 지원
        ".docx": extract_text_from_docx,
        ".hwp": extract_text_from_hwp,
        ".hwpx": extract_text_from_hwp,  # hwpx는 다른 처리 필요할 수 있음
        ".ppt": extract_text_from_pptx,
        ".pptx": extract_text_from_pptx,
        ".xls": extract_text_from_xlsx,
        ".xlsx": extract_text_from_xlsx,
    }
    
    extractor = extractors.get(ext)
    if extractor:
        return extractor(filepath)
    else:
        print(f"[경고] 지원하지 않는 파일 형식: {ext}")
        return ""


# ======================
# 분류기
# ======================

class DocumentClassifier:
    def __init__(self, model_dir: str, device: str = None):
        self.device = torch.device(device if device else ("cuda" if torch.cuda.is_available() else "cpu"))
        
        # 라벨 매핑 로드
        label_map_path = os.path.join(model_dir, "label_mapping.json")
        with open(label_map_path, "r", encoding="utf-8") as f:
            mapping = json.load(f)
            self.label2id = mapping["label2id"]
            self.id2label = {int(k): v for k, v in mapping["id2label"].items()}
        
        # 토크나이저 로드
        self.tokenizer = AutoTokenizer.from_pretrained(model_dir)
        
        # 모델 로드
        num_classes = len(self.label2id)
        self.model = KoBERTClassifier(num_classes=num_classes)
        model_path = os.path.join(model_dir, "best_model.pt")
        self.model.load_state_dict(torch.load(model_path, map_location=self.device))
        self.model.to(self.device)
        self.model.eval()
        
        print(f"[INFO] 모델 로드 완료: {model_dir}")
        print(f"[INFO] 분류 가능한 문서 종류: {list(self.label2id.keys())}")
    
    def predict(self, text: str, max_length: int = 512) -> tuple:
        """텍스트를 분류하고 (라벨, 확률) 반환"""
        if not text.strip():
            return None, 0.0
        
        encoding = self.tokenizer(
            text,
            truncation=True,
            padding="max_length",
            max_length=max_length,
            return_tensors="pt",
        )
        
        input_ids = encoding["input_ids"].to(self.device)
        attention_mask = encoding["attention_mask"].to(self.device)
        
        with torch.no_grad():
            logits = self.model(input_ids, attention_mask)
            probs = torch.softmax(logits, dim=1)
            confidence, pred_id = torch.max(probs, dim=1)
        
        label = self.id2label[pred_id.item()]
        return label, confidence.item()


# ======================
# 메인 로직
# ======================

def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--model_dir",
        type=str,
        default="./saved_model",
        help="학습된 모델이 저장된 디렉토리 (기본: ./saved_model)",
    )
    parser.add_argument(
        "--input_dir",
        type=str,
        default="./input",
        help="분류할 문서들이 있는 입력 폴더 (기본: ./input)",
    )
    parser.add_argument(
        "--output_dir",
        type=str,
        default="./output",
        help="분류 결과를 저장할 출력 폴더 (기본: ./output)",
    )
    parser.add_argument(
        "--min_confidence",
        type=float,
        default=0.5,
        help="최소 신뢰도. 이보다 낮으면 '미분류' 폴더로 이동 (기본: 0.5)",
    )
    parser.add_argument(
        "--copy",
        action="store_true",
        help="파일을 이동 대신 복사 (기본: 이동)",
    )
    return parser.parse_args()


def main():
    args = parse_args()
    
    input_dir = Path(args.input_dir)
    output_dir = Path(args.output_dir)
    
    if not input_dir.exists():
        print(f"[에러] 입력 폴더가 없습니다: {input_dir}")
        return
    
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # 분류기 로드
    classifier = DocumentClassifier(args.model_dir)
    
    # 지원하는 확장자
    supported_exts = {".txt", ".pdf", ".doc", ".docx", ".hwp", ".hwpx", ".ppt", ".pptx", ".xls", ".xlsx"}
    
    # 통계
    stats = {"total": 0, "success": 0, "failed": 0, "low_confidence": 0}
    
    # 입력 폴더의 파일만 처리 (하위 폴더는 건드리지 않음)
    all_files = [f for f in input_dir.iterdir() if f.is_file() and f.suffix.lower() in supported_exts]
    
    print(f"[INFO] 처리할 파일: {len(all_files)}개")
    
    for filepath in all_files:
        ext = filepath.suffix.lower()
        
        stats["total"] += 1
        print(f"\n[처리중] {filepath.name}")
        
        # 텍스트 추출
        text = extract_text(filepath)
        if not text.strip():
            print(f"  -> 텍스트 추출 실패")
            dest_folder = output_dir / "추출실패"
            dest_folder.mkdir(exist_ok=True)
            dest_path = dest_folder / filepath.name
            if args.copy:
                shutil.copy2(filepath, dest_path)
            else:
                shutil.move(str(filepath), str(dest_path))
            stats["failed"] += 1
            continue
        
        # 분류
        label, confidence = classifier.predict(text)
        print(f"  -> 분류: {label} (신뢰도: {confidence:.2%})")
        
        # 신뢰도가 낮으면 미분류로
        if confidence < args.min_confidence:
            label = "미분류"
            stats["low_confidence"] += 1
        else:
            stats["success"] += 1
        
        # 파일 이동/복사
        dest_folder = output_dir / label
        dest_folder.mkdir(exist_ok=True)
        dest_path = dest_folder / filepath.name
        
        # 중복 파일명 처리
        counter = 1
        while dest_path.exists():
            stem = filepath.stem
            dest_path = dest_folder / f"{stem}_{counter}{ext}"
            counter += 1
        
        if args.copy:
            shutil.copy2(filepath, dest_path)
            print(f"  -> 복사: {dest_path}")
        else:
            shutil.move(str(filepath), str(dest_path))
            print(f"  -> 이동: {dest_path}")
    
    # 결과 요약
    print("\n" + "=" * 50)
    print("[완료] 분류 결과")
    print(f"  - 총 파일: {stats['total']}개")
    print(f"  - 분류 성공: {stats['success']}개")
    print(f"  - 신뢰도 부족 (미분류): {stats['low_confidence']}개")
    print(f"  - 텍스트 추출 실패: {stats['failed']}개")


if __name__ == "__main__":
    main()

