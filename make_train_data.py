"""
input 폴더의 문서들을 자동으로 train.json으로 변환하는 스크립트

폴더 구조:
input/
├── 기안서/
│   ├── 문서1.pdf
│   └── 문서2.docx
├── 견적서/
│   ├── 견적서1.pdf
│   └── 견적서2.hwp
├── 계약서/
│   └── ...
└── ...

→ 폴더명이 라벨(label)이 됩니다!

사용 예시:
python make_train_data.py --input_dir ./input --output_file ./data/train.json
"""

import os
import json
import argparse
from pathlib import Path
from typing import List, Dict


# ======================
# 문서에서 텍스트 추출
# ======================

def extract_text_from_txt(filepath: Path) -> str:
    """TXT 파일에서 텍스트 추출"""
    encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
    for enc in encodings:
        try:
            with open(filepath, 'r', encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    return ""


def extract_text_from_pdf(filepath: Path) -> str:
    """PDF 파일에서 텍스트 추출"""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n".join(text_parts)
    except ImportError:
        print("[경고] pdfplumber 설치 필요: pip install pdfplumber")
        return ""
    except Exception as e:
        print(f"[경고] PDF 읽기 실패: {filepath} ({e})")
        return ""


def extract_text_from_docx(filepath: Path) -> str:
    """DOCX 파일에서 텍스트 추출"""
    try:
        from docx import Document
        doc = Document(filepath)
        text_parts = [para.text for para in doc.paragraphs]
        return "\n".join(text_parts)
    except ImportError:
        print("[경고] python-docx 설치 필요: pip install python-docx")
        return ""
    except Exception as e:
        print(f"[경고] DOCX 읽기 실패: {filepath} ({e})")
        return ""


def extract_text_from_hwp(filepath: Path) -> str:
    """HWP 파일에서 텍스트 추출"""
    try:
        import olefile
        with olefile.OleFileIO(filepath) as ole:
            if ole.exists("PrvText"):
                encoded_text = ole.openstream("PrvText").read()
                return encoded_text.decode('utf-16', errors='ignore')
        return ""
    except ImportError:
        print("[경고] olefile 설치 필요: pip install olefile")
        return ""
    except Exception as e:
        print(f"[경고] HWP 읽기 실패: {filepath} ({e})")
        return ""


def extract_text_from_pptx(filepath: Path) -> str:
    """PPTX 파일에서 텍스트 추출"""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except ImportError:
        print("[경고] python-pptx 설치 필요: pip install python-pptx")
        return ""
    except Exception as e:
        print(f"[경고] PPTX 읽기 실패: {filepath} ({e})")
        return ""


def extract_text_from_xlsx(filepath: Path) -> str:
    """XLSX 파일에서 텍스트 추출"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text_parts.append(str(cell.value))
        return " ".join(text_parts)
    except ImportError:
        print("[경고] openpyxl 설치 필요: pip install openpyxl")
        return ""
    except Exception as e:
        print(f"[경고] XLSX 읽기 실패: {filepath} ({e})")
        return ""


def extract_text(filepath: Path) -> str:
    """파일 확장자에 따라 텍스트 추출"""
    ext = filepath.suffix.lower()
    
    extractors = {
        ".txt": extract_text_from_txt,
        ".pdf": extract_text_from_pdf,
        ".doc": extract_text_from_docx,
        ".docx": extract_text_from_docx,
        ".hwp": extract_text_from_hwp,
        ".hwpx": extract_text_from_hwp,
        ".ppt": extract_text_from_pptx,
        ".pptx": extract_text_from_pptx,
        ".xls": extract_text_from_xlsx,
        ".xlsx": extract_text_from_xlsx,
    }
    
    extractor = extractors.get(ext)
    if extractor:
        return extractor(filepath)
    else:
        print(f"[경고] 지원하지 않는 파일 형식: {ext}")
        return ""


def clean_text(text: str, max_length: int = 2000) -> str:
    """텍스트 정리 (공백 정리, 길이 제한)"""
    # 연속 공백/줄바꿈 정리
    text = " ".join(text.split())
    # 너무 길면 자르기 (KoBERT max_length 고려)
    if len(text) > max_length:
        text = text[:max_length]
    return text


# ======================
# 메인 로직
# ======================

def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--input_dir",
        type=str,
        default="./input",
        help="문서가 있는 입력 폴더 (하위 폴더명 = 라벨)",
    )
    parser.add_argument(
        "--output_file",
        type=str,
        default="./data/train.json",
        help="생성할 train.json 파일 경로",
    )
    parser.add_argument(
        "--max_length",
        type=int,
        default=2000,
        help="텍스트 최대 길이 (기본: 2000자)",
    )
    parser.add_argument(
        "--min_length",
        type=int,
        default=50,
        help="텍스트 최소 길이 (이보다 짧으면 제외, 기본: 50자)",
    )
    return parser.parse_args()


def main():
    args = parse_args()
    
    input_dir = Path(args.input_dir)
    output_file = Path(args.output_file)
    
    if not input_dir.exists():
        print(f"[에러] 입력 폴더가 없습니다: {input_dir}")
        print("\n폴더 구조를 만들어주세요:")
        print("  input/")
        print("  ├── 기안서/")
        print("  │   └── (PDF/DOCX/HWP 파일들)")
        print("  ├── 견적서/")
        print("  │   └── (PDF/DOCX/HWP 파일들)")
        print("  └── ...")
        return
    
    # 지원하는 확장자
    supported_exts = {".txt", ".pdf", ".doc", ".docx", ".hwp", ".hwpx", ".ppt", ".pptx", ".xls", ".xlsx"}
    
    # 결과 저장
    train_data: List[Dict] = []
    stats = {}
    
    # 하위 폴더 탐색 (폴더명 = 라벨)
    for label_dir in input_dir.iterdir():
        if not label_dir.is_dir():
            continue
        
        label = label_dir.name
        print(f"\n[라벨] {label}")
        stats[label] = {"total": 0, "success": 0, "failed": 0}
        
        # 해당 폴더의 모든 파일 처리
        for filepath in label_dir.iterdir():
            if not filepath.is_file():
                continue
            
            ext = filepath.suffix.lower()
            if ext not in supported_exts:
                continue
            
            stats[label]["total"] += 1
            print(f"  - {filepath.name}", end=" ")
            
            # 텍스트 추출
            text = extract_text(filepath)
            text = clean_text(text, args.max_length)
            
            if len(text) < args.min_length:
                print(f"→ 텍스트 부족 ({len(text)}자)")
                stats[label]["failed"] += 1
                continue
            
            train_data.append({
                "text": text,
                "label": label
            })
            stats[label]["success"] += 1
            print(f"→ OK ({len(text)}자)")
    
    if not train_data:
        print("\n[에러] 추출된 데이터가 없습니다!")
        return
    
    # 출력 폴더 생성
    output_file.parent.mkdir(parents=True, exist_ok=True)
    
    # JSON 저장
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(train_data, f, ensure_ascii=False, indent=2)
    
    # 결과 요약
    print("\n" + "=" * 50)
    print("[완료] 학습 데이터 생성 결과")
    print(f"  - 출력 파일: {output_file}")
    print(f"  - 총 데이터 수: {len(train_data)}개")
    print("\n  라벨별 통계:")
    for label, stat in stats.items():
        print(f"    - {label}: {stat['success']}개 성공 / {stat['total']}개 중")
    
    print(f"\n이제 학습을 실행하세요:")
    print(f"  python train.py --train_file {output_file} --output_dir ./saved_model")


if __name__ == "__main__":
    main()

